import shutil
import os

# Auto delete ChromaDB folder if it exists
chroma_path = "./chroma"  # or "./.chromadb" if that's your folder
if os.path.exists(chroma_path):
    shutil.rmtree(chroma_path)
    print(f"Deleted old ChromaDB folder at {chroma_path}")


import streamlit as st
import logging
import os
import tempfile
import shutil
import pdfplumber
import ollama
import torch
import soundfile as sf
import warnings
warnings.filterwarnings("ignore")

from langchain_community.document_loaders import UnstructuredPDFLoader  
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.embeddings import OllamaEmbeddings
from langchain_ollama import ChatOllama
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain.prompts import ChatPromptTemplate, PromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_core.runnables import RunnablePassthrough
from langchain.retrievers.multi_query import MultiQueryRetriever
from transformers import AutoProcessor, AutoModelForSpeechSeq2Seq
from typing import List, Tuple, Dict, Any, Optional
from pptx import Presentation

# Initialize Whisper model and processor
processor = AutoProcessor.from_pretrained("openai/whisper-small")
model = AutoModelForSpeechSeq2Seq.from_pretrained("openai/whisper-small")
import numpy as np
import torch
import soundfile as sf
from transformers import AutoProcessor, AutoModelForSpeechSeq2Seq

# Initialize Whisper model and processor
processor = AutoProcessor.from_pretrained("openai/whisper-small")
model = AutoModelForSpeechSeq2Seq.from_pretrained("openai/whisper-small")

def transcribe_audio_with_whisper(audio_file_path: str, chunk_duration=30) -> str:
    """
    Transcribe audio using OpenAI Whisper, splitting into smaller chunks to manage memory.
    """
    # Read the audio file
    audio_input, sample_rate = sf.read(audio_file_path)
    
    # Calculate the number of samples for the chunk duration
    chunk_size = chunk_duration * sample_rate  # in samples
    num_chunks = len(audio_input) // chunk_size + 1  # number of chunks
    
    # Prepare to store transcriptions
    transcriptions = []
    
    for i in range(num_chunks):
        # Extract a chunk of audio
        start_sample = i * chunk_size
        end_sample = min((i + 1) * chunk_size, len(audio_input))
        audio_chunk = audio_input[start_sample:end_sample]
        
        # Process the chunk
        inputs = processor(audio_chunk, return_tensors="pt", sampling_rate=sample_rate)
        
        # Generate transcription
        with torch.no_grad():
            generated_ids = model.generate(inputs["input_features"])
        
        # Decode the transcription and add it to the list
        transcription = processor.batch_decode(generated_ids, skip_special_tokens=True)[0]
        transcriptions.append(transcription)
    
    # Join all transcriptions from chunks
    return " ".join(transcriptions).strip()


# Streamlit page configuration
st.set_page_config(
    page_title="Interactive File Query GPT",
    page_icon="🎈",
    layout="wide",
    initial_sidebar_state="collapsed",
)

# Logging configuration
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
)

logger = logging.getLogger(__name__)

@st.cache_resource(show_spinner=True)
def extract_model_names_from_serializable(model_names_list: List[str]) -> Tuple[str, ...]:
    return tuple(model_names_list)

def create_vector_db_from_text(text: str) -> Chroma:
    """
    Create a vector database from transcribed text.
    """
    logger.info(f"Creating vector DB from transcribed text")
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=7500, chunk_overlap=100)
    chunks = text_splitter.split_text(text)
    logger.info("Text split into chunks")

    # Assuming you use streamlit to get the file and then process embeddings
   
           
    embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")

    # Proceed with further processing...


    vector_db = Chroma.from_texts(texts=chunks, embedding=embeddings, collection_name="myRAG")
    logger.info("Vector DB created")
    return vector_db
def create_vector_db(file_upload) -> Chroma:
    """
    Create a vector database from an uploaded PDF file with fallback OCR if needed.
    """
    logger.info(f"Creating vector DB from file upload: {file_upload.name}")
    temp_dir = tempfile.mkdtemp()

    path = os.path.join(temp_dir, file_upload.name)
    with open(path, "wb") as f:
        f.write(file_upload.getvalue())
        logger.info(f"File saved to temporary path: {path}")

    try:
        # PDF extraction using pdfplumber
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            data = []
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    data.append(Document(page_content=text))
            
            if not data:
                st.error("❌ Failed to extract any content from the uploaded PDF using pdfplumber. Try another file or check its format.")
                shutil.rmtree(temp_dir)
                return None

        logger.info("PDF text extraction succeeded with pdfplumber.")

    except Exception as e:
        logger.warning(f"PDF text extraction failed: {e}, trying OCR fallback...")

        try:
            from pdf2image import convert_from_path
            import pytesseract
            from langchain_core.documents import Document

            images = convert_from_path(path)
            ocr_text = "\n".join([pytesseract.image_to_string(img) for img in images])
            data = [Document(page_content=ocr_text)]
            logger.info("OCR-based text extraction succeeded.")
        except Exception as ocr_error:
            logger.error(f"OCR fallback also failed: {ocr_error}")
            data = []

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=7500, chunk_overlap=100)
    chunks = text_splitter.split_documents(data)
    logger.info("Document split into chunks")

    embeddings = OllamaEmbeddings(model="nomic-embed-text")
    vector_db = Chroma.from_documents(
        documents=chunks, embedding=embeddings, collection_name="myRAG"
    )
    logger.info("Vector DB created")

    shutil.rmtree(temp_dir)
    logger.info(f"Temporary directory {temp_dir} removed")
    return vector_db


def process_question(question: str, vector_db: Chroma, selected_model: str) -> str:
    """
    Process a user question using the vector database and selected language model.
    """
    logger.info(f"""Processing question: {question} using model: {selected_model}""")
    llm = ChatOllama(model=selected_model, temperature=0)
    QUERY_PROMPT = PromptTemplate(
        input_variables=["question"],
        template="""You are an AI language model assistant. Your task is to generate 3
        different versions of the given user question to retrieve relevant documents from
        a vector database. Add the phrase "According to the document" in each question.
        By generating multiple perspectives on the user question, your
        goal is to help the user overcome some of the limitations of the distance-based
        similarity search. Provide these alternative questions separated by newlines.
        Original question: {question}""",
    )

    retriever = MultiQueryRetriever.from_llm(
        vector_db.as_retriever(), llm, prompt=QUERY_PROMPT
    )

    template = """Answer the question based ONLY on the following context:
    {context}
    Question: {question}
    If you don't know the answer, just say that you don't know, don't try to make up an answer.
    Only provide the answer from the {context}, nothing else.
    Add snippets of the context you used to answer the question.
    Only answer the question if the document contains related matter.
    """

    prompt = ChatPromptTemplate.from_template(template)

    chain = (
        {"context": retriever, "question": RunnablePassthrough()}
        | prompt
        | llm
        | StrOutputParser()
    )

    response = chain.invoke(question)
    logger.info("Question processed and response generated")
    return response

@st.cache_data
def extract_all_pages_as_images(file_upload) -> List[Any]:
    """

    Extract all pages from a PDF file as images.
    """
    logger.info(f"""Extracting all pages as images from file: {file_upload.name}""")
    pdf_pages = []
    with pdfplumber.open(file_upload) as pdf:
        pdf_pages = [page.to_image().original for page in pdf.pages]
    logger.info("PDF pages extracted as images")
    return pdf_pages

def delete_vector_db(vector_db: Optional[Chroma]) -> None:
    """
    Delete the vector database and clear related session state.
    """
    logger.info("Deleting vector DB")
    if vector_db is not None:
        
        vector_db.delete_collection()
        st.session_state.pop("pdf_pages", None)
        st.session_state.pop("file_upload", None)
        st.session_state.pop("vector_db", None)
        st.success("Collection and temporary files deleted successfully.")
        logger.info("Vector DB and related session state cleared")
        st.rerun()
    else:
        st.error("No vector database found to delete.")
        logger.warning("Attempted to delete vector DB, but none was found")

def transcribe_audio_with_whisper(audio_file_path: str) -> str:
    
    """
    Transcribe audio using OpenAI Whisper.
    """
    audio_input, _ = sf.read(audio_file_path)
    inputs = processor(audio_input, return_tensors="pt", sampling_rate=16000)
    with torch.no_grad():
        generated_ids = model.generate(inputs["input_features"])
    transcription = processor.batch_decode(generated_ids, skip_special_tokens=True)[0]
    return transcription.strip()
def extract_text_from_ppt(ppt_file_path: str) -> str:
    """
    Extract text from a PPT file.
    """
    prs = Presentation(ppt_file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    return text.strip()

def main() -> None:
    """
    Main function to run the Streamlit application.

    This function sets up the user interface, handles file uploads,
    processes user queries, and displays results.
    """
    st.subheader("Interactive File Query GPT", divider="gray", anchor=False)

    # Fetch models from Ollama
    models_info = ollama.list()
    models_raw = models_info.get("models", [])


    try:
        models_raw = models_info.get("models", [])
        model_names_list = [
        model.model
        for model in models_raw
    ]
    except Exception as e:
        logger.error(f"Error parsing model names: {e}")
        model_names_list = []

    available_models = extract_model_names_from_serializable(model_names_list)
    col1, col2 = st.columns([1.5, 2])

    selected_model = None  # Fallback
    if available_models:
        selected_model = col2.selectbox(
            "Pick a model available locally on your system ↓", available_models
        )
    else:
        st.warning("⚠️ No Ollama models detected. Make sure Ollama is running and models are pulled.")

    # Remaining logic (file upload, chat input, etc.) should follow here...


# Optional: Print the result to check the output
    logger.info(f"Ollama model info raw response: {models_info}")


    

    col1, col2 = st.columns([1.5, 2])

    if "messages" not in st.session_state:
        st.session_state["messages"] = []

    if "vector_db" not in st.session_state:
        st.session_state["vector_db"] = None

   
        
    file_upload = col1.file_uploader(
        "Upload a PDF file ↓", type="pdf", accept_multiple_files=False
    )

    audio_upload = col1.file_uploader(
        "Upload an audio file ↓", type=["wav", "mp3"], accept_multiple_files=False
    )

    ppt_upload = col1.file_uploader(
        "Upload a PPT file ↓", type=["ppt", "pptx"], accept_multiple_files=False
    )

    txt_upload = col1.file_uploader(
        "Upload a TXT file ↓", type="txt", accept_multiple_files=False
    )
     
    if file_upload:
        st.session_state["file_upload"] = file_upload
        if st.session_state["vector_db"] is None:
            with st.spinner("📚 Processing PDF..."):
                 st.session_state["vector_db"] = create_vector_db(file_upload)       
        pdf_pages = extract_all_pages_as_images(file_upload)
        st.session_state["pdf_pages"] = pdf_pages

        zoom_level = col1.slider(
            "Zoom Level", min_value=100, max_value=1000, value=700, step=50
        )

        with col1:
            with st.container(height=410, border=True):
                for page_image in pdf_pages:
                    st.image(page_image, width=zoom_level)

    if ppt_upload:
        with st.spinner("Processing PPT file..."):
            temp_ppt_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
            temp_ppt_file.write(ppt_upload.read())
            temp_ppt_file.flush()
            temp_ppt_file.close()

            # Extract text from the PPT file
            text = extract_text_from_ppt(temp_ppt_file.name)
            os.remove(temp_ppt_file.name)
            st.success("PPT processing completed.")

            # Store the extracted text in the vector database
            if st.session_state["vector_db"] is None:
                st.session_state["vector_db"] = create_vector_db_from_text(text)
            else:
                combined_text = text #+ " " + "\n".join([doc.page_content for doc in st.session_state["vector_db"].get_all_documents()])
                st.session_state["vector_db"] = create_vector_db_from_text(combined_text)

    if txt_upload:
        with st.spinner("Processing TXT file..."):
            temp_txt_file = tempfile.NamedTemporaryFile(delete=False, suffix=".txt")
            temp_txt_file.write(txt_upload.read())
            temp_txt_file.flush()
            temp_txt_file.close()

            with open(temp_txt_file.name, 'r') as file:
                text = file.read()

            os.remove(temp_txt_file.name)
            st.success("TXT processing completed.")

            # Store the extracted text in the vector database
            if st.session_state["vector_db"] is None:
                st.session_state["vector_db"] = create_vector_db_from_text(text)
            else:
                combined_text = text #+ " " + "\n".join([doc.page_content for doc in st.session_state["vector_db"].get_all_documents()])
                st.session_state["vector_db"] = create_vector_db_from_text(combined_text)

    delete_collection = col1.button("⚠️ Delete collection", type="secondary")

    if delete_collection:
        delete_vector_db(st.session_state["vector_db"])

    with col2:
        message_container = st.container(height=500, border=True)

        for message in st.session_state["messages"]:
            avatar = "🤖" if message["role"] == "assistant" else "😎"
            with message_container.chat_message(message["role"], avatar=avatar):
                st.markdown(message["content"])

        if prompt := st.chat_input("Enter a prompt here..."):
            try:
                st.session_state["messages"].append({"role": "user", "content": prompt})
                message_container.chat_message("user", avatar="😎").markdown(prompt)

                with message_container.chat_message("assistant", avatar="🤖"):
                    with st.spinner(":green[processing...]"):
                        if st.session_state["vector_db"] is not None:
                            response = process_question(
                                prompt, st.session_state["vector_db"], selected_model
                            )
                            st.markdown(response)
                        else:
                            st.warning("Please upload a file first.")

                if st.session_state["vector_db"] is not None:
                    st.session_state["messages"].append(
                        {"role": "assistant", "content": response}
                    )

            except Exception as e:
                st.error(e, icon="⛔️")
                logger.error(f"Error processing prompt: {e}")
        else:
            if st.session_state["vector_db"] is None:
                st.warning("Upload a file to begin chat...")

    if audio_upload:
        with st.spinner("Transcribing audio..."):
            temp_audio_file = tempfile.NamedTemporaryFile(delete=False, suffix=".wav")
            temp_audio_file.write(audio_upload.read())
            temp_audio_file.flush()
            temp_audio_file.close()
            transcription = transcribe_audio_with_whisper(temp_audio_file.name)
            os.remove(temp_audio_file.name)
            st.success("Audio transcription completed.")
            st.text_area("Transcription:", transcription, height=200)

            # Store the transcription in the vector database
            if st.session_state["vector_db"] is None:
                st.session_state["vector_db"] = create_vector_db_from_text(transcription)
            else:
                combined_text = transcription #+ " " + "\n".join([doc.page_content for doc in st.session_state["vector_db"].get_all_documents()])
                st.session_state["vector_db"] = create_vector_db_from_text(combined_text)

if __name__ == "__main__":
    main()